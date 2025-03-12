import io
import pandas as pd
import docx
import pdfplumber
from pptx import Presentation
from openpyxl import load_workbook

def extract_text(file_name, file_bytes):
    """Extracts text from various document formats using in-memory file content."""
    file_extension = file_name.split(".")[-1].lower()
    
    if file_extension == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif file_extension == "docx":
        return extract_text_from_docx(file_bytes)
    elif file_extension == "txt":
        return extract_text_from_txt(file_bytes)
    elif file_extension in ["xlsx", "csv"]:
        return extract_text_from_excel_csv(file_bytes, file_extension)
    elif file_extension == "pptx":
        return extract_text_from_pptx(file_bytes)
    else:
        return "Unsupported file format."

def extract_text_from_pdf(file_bytes):
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        return "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())

def extract_text_from_docx(file_bytes):
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(file_bytes):
    return file_bytes.decode("utf-8")

def extract_text_from_excel_csv(file_bytes, file_extension):
    if file_extension == "csv":
        df = pd.read_csv(io.BytesIO(file_bytes))
    else:
        df = pd.read_excel(io.BytesIO(file_bytes), engine='openpyxl')
    return df.to_string()

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
